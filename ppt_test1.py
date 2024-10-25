from pptx import Presentation

# プレゼンテーションオブジェクトを作成
prs = Presentation()

# タイトルスライドを追加
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)

# タイトルとサブタイトルを設定
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Pythonで作成したプレゼンテーション"
subtitle.text = "python-pptxライブラリを使用"

# 新しいスライドを追加
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# スライドのタイトルとコンテンツを設定
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "新しいスライド"
content.text = "ここにコンテンツを追加します。"

# プレゼンテーションを保存
prs.save('test_presentation.pptx')
