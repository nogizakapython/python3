#################################################
#   Excelからpowerpointの表にデータコピーテスト
#################################################

# ライブラリをインポート
import openpyxl
from pptx import Presentation
from pptx.util import Inches
# Excel表データ用配列
excel_data = []


# Excelの表から、PowerPointの表に転記するデータを取り出し、配列に格納する。
def excel_data_copy():
    # Excelファイルを開く
    excel_file = '調査結果.xlsx'
    sheet_name = '調査結果'

    # Excelファイルが開いていれば例外処理に遷移し、ファイルを閉じるメッセージを出力して終了する。
    try:
        wb = openpyxl.load_workbook(excel_file)
        ws = wb[sheet_name]
    except:
        print(f"{excel_file}が開いています。閉じてください")    



    # Excelデータ開始変数
    start_row = 4
    # Excelデータ終了行変数
    end_row = 0
    # カウント変数
    i = start_row
    while True:
        v1 = ws.cell(row=i,column=3).value
        if v1 == None:
            break
        else:
            i += 1
    
    # 最終行を取得
    end_row = i-1
    for row in ws.iter_rows(values_only=True,min_row=start_row,max_row=end_row,min_col=3):
        excel_data.append(row)

def powerpoint_table_copy():
    # PowerPointプレゼンテーションを作成
    ppt_file = "プレゼンテスト6.pptx"
    ppt = Presentation(ppt_file)
    # スライドの２ページ目に追加
    slide = ppt.slides[1]
    # 表を追加
    rows = 2
    cols = 6
    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(8.0)
    height = Inches(1.5)
    # テーブルの定義
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for j in range(rows):
        array1 = excel_data[j]
        leng1 = len(array1)
        for k in range(leng1):
            data2 = array1[k]
            # 文字列型にして表に貼り付け
            table.cell(j, k).text = str(data2)
        

    # PowerPointファイルを保存
    try:
        ppt.save(ppt_file)
    except PermissionError:
        print(f"{ppt_file}が開いています。閉じてください") 

# メイン関数
def main():
    excel_data_copy()
    powerpoint_table_copy()

# 主処理
if __name__ == "__main__":
    main()            