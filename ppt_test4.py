# ライブラリのインポート
from pptx import Presentation
from pptx.util import Inches

# 新しいプレゼンテーションを作成
prs = Presentation()

# スライドを追加
slide_layout = prs.slide_layouts[5]  # 空のスライドレイアウトを選択
slide = prs.slides.add_slide(slide_layout)

# スライドのタイトルとコンテンツを設定
title = slide.shapes.title
title.text = "縄跳びの回数"
# 表を追加
rows = 2
cols = 2
left = Inches(2.0)
top = Inches(2.0)
width = Inches(4.0)
height = Inches(1.0)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# 表のセルにデータを追加
table.cell(0, 0).text = '日付'
table.cell(0, 1).text = '回数'
table.cell(1, 0).text = '2024/11/4'
table.cell(1, 1).text = '5'

# プレゼンテーションを保存
try:
    prs.save('プレゼンテスト3.pptx')
except PermissionError:
    print("ファイルを閉じてください")
      