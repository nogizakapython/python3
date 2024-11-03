# インポートしないとエラーが発生する
# import collections.abc
#  プレゼンテーションを作成
from pptx import Presentation 
# インチ
# from pptx.util import Inches  
# from pptx.enum.text import PP_ALIGN  # 中央揃えにする用
# from pptx.util import Cm, Pt # センチ、ポイント


def slide_count1(obj1):
    count = len(obj1)
    return count


# プレゼンテーションオブジェクトを作成
prs = Presentation()

# プレゼンテーションファイルを開く
filepath = "プレゼンテスト1.pptx"

# 例外処理
try:
    ppt = Presentation(filepath)
except:
    print(f"{filepath}がありません！")
    exit()    


# スライド数を取得する
slide_count = slide_count1(ppt.slides)
print(slide_count)

# プレゼンテーションを保存
prs.save('test_presentation.pptx')
